"""Structure-aware document parsing service."""
import os
import re
from pathlib import Path
from typing import Iterable

from app.core.logging import get_logger
from app.services.document_block import NoExtractableTextError, ParsedBlock

logger = get_logger(__name__)


class DocumentParser:
    """Parse supported files into structure-aware blocks."""

    SUPPORTED_TYPES = {
        "pdf": ["pdf"],
        "word": ["docx"],
        "excel": ["xlsx"],
        "pptx": ["pptx"],
        "text": ["txt", "md", "csv", "json", "xml", "html", "htm"],
    }
    _HEADING_STYLE = re.compile(r"^(?:heading|标题)\s*([1-9])", re.IGNORECASE)
    _MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.+?)\s*$")

    def get_file_type(self, filename: str) -> str:
        ext = Path(filename).suffix.lower().lstrip(".")
        for file_type, extensions in self.SUPPORTED_TYPES.items():
            if ext in extensions:
                return file_type
        return "unknown"

    def parse(self, file_path: str, filename: str | None = None) -> tuple[list[ParsedBlock], dict]:
        filename = filename or os.path.basename(file_path)
        ext = Path(filename).suffix.lower().lstrip(".")
        file_type = self.get_file_type(filename)

        logger.info("Parsing file: %s (type=%s)", filename, file_type)
        if file_type == "pdf":
            blocks = self._parse_pdf(file_path)
        elif file_type == "word":
            blocks = self._parse_word(file_path)
        elif file_type == "excel":
            blocks = self._parse_excel(file_path)
        elif file_type == "pptx":
            blocks = self._parse_pptx(file_path)
        elif file_type == "text":
            blocks = self._parse_text(file_path, ext)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        blocks = [block for block in blocks if block.normalized_text()]
        if not blocks:
            raise NoExtractableTextError()

        document_title = Path(filename).stem.strip() or filename
        for block in blocks:
            block.metadata.setdefault("document_title", document_title)
            block.metadata.setdefault("source", filename)
            block.metadata.setdefault("block_type", block.block_type)
            if block.heading_path:
                block.metadata.setdefault("heading_path", " > ".join(block.heading_path))

        total_chars = sum(len(block.page_content) for block in blocks)
        meta = {
            "filename": filename,
            "file_type": file_type,
            "ext": ext,
            "block_count": len(blocks),
            "page_count": len(blocks),
            "total_chars": total_chars,
        }
        logger.info("Parsed %d blocks, %d chars from %s", len(blocks), total_chars, filename)
        return blocks, meta

    def _parse_pdf(self, file_path: str) -> list[ParsedBlock]:
        import fitz

        blocks: list[ParsedBlock] = []
        pdf_document = fitz.open(file_path)
        try:
            for page_number, page in enumerate(pdf_document, 1):
                text = (page.get_text("text") or "").strip()
                if text:
                    blocks.append(ParsedBlock(
                        page_content=text,
                        block_type="page",
                        metadata={"page": page_number},
                    ))
        finally:
            pdf_document.close()
        return blocks

    def _parse_word(self, file_path: str) -> list[ParsedBlock]:
        from docx import Document as DocxDocument
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = DocxDocument(file_path)
        blocks: list[ParsedBlock] = []
        headings: list[tuple[int, str]] = []
        paragraphs: list[str] = []
        pending_heading = ""

        def heading_path() -> tuple[str, ...]:
            return tuple(title for _, title in headings)

        def flush_paragraphs() -> None:
            if paragraphs:
                blocks.append(ParsedBlock(
                    page_content="\n\n".join(paragraphs),
                    block_type="prose",
                    heading_path=heading_path(),
                ))
                paragraphs.clear()

        def flush_standalone_heading() -> None:
            nonlocal pending_heading
            if pending_heading:
                blocks.append(ParsedBlock(
                    page_content=pending_heading,
                    block_type="heading",
                    heading_path=heading_path(),
                ))
                pending_heading = ""

        for element in document.element.body.iterchildren():
            if isinstance(element, CT_P):
                paragraph = Paragraph(element, document)
                text = paragraph.text.strip()
                if not text:
                    continue
                style_name = paragraph.style.name if paragraph.style is not None else ""
                heading_match = self._HEADING_STYLE.match(style_name or "")
                if heading_match:
                    flush_paragraphs()
                    flush_standalone_heading()
                    level = int(heading_match.group(1))
                    headings[:] = [(heading_level, title) for heading_level, title in headings
                                   if heading_level < level]
                    headings.append((level, text))
                    pending_heading = text
                else:
                    paragraphs.append(text)
                    pending_heading = ""
            elif isinstance(element, CT_Tbl):
                flush_paragraphs()
                table = Table(element, document)
                rows = [[self._clean_table_cell(cell.text) for cell in row.cells] for row in table.rows]
                markdown = self._rows_to_markdown(rows)
                if markdown:
                    blocks.append(ParsedBlock(
                        page_content=markdown,
                        block_type="table",
                        heading_path=heading_path(),
                        metadata={"table_format": "markdown"},
                    ))
                    pending_heading = ""

        flush_paragraphs()
        flush_standalone_heading()
        return blocks

    def _parse_excel(self, file_path: str) -> list[ParsedBlock]:
        import openpyxl

        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        blocks: list[ParsedBlock] = []
        try:
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    values = [self._clean_table_cell(cell) for cell in row]
                    if any(values):
                        rows.append(values)
                markdown = self._rows_to_markdown(rows)
                if markdown:
                    blocks.append(ParsedBlock(
                        page_content=markdown,
                        block_type="table",
                        heading_path=(sheet_name,),
                        metadata={
                            "sheet": sheet_name,
                            "table_format": "markdown",
                            "header_inferred": True,
                        },
                    ))
        finally:
            workbook.close()
        return blocks

    def _parse_pptx(self, file_path: str) -> list[ParsedBlock]:
        from pptx import Presentation

        presentation = Presentation(file_path)
        blocks: list[ParsedBlock] = []
        for slide_number, slide in enumerate(presentation.slides, 1):
            title = ""
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
            texts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text and text != title:
                    texts.append(text)
            content = "\n\n".join(texts) or title
            if content:
                blocks.append(ParsedBlock(
                    page_content=content,
                    block_type="slide",
                    heading_path=(title,) if title else (),
                    metadata={"slide": slide_number},
                ))
        return blocks

    def _parse_text(self, file_path: str, extension: str) -> list[ParsedBlock]:
        import chardet

        with open(file_path, "rb") as file:
            raw = file.read()
        encoding = chardet.detect(raw).get("encoding", "utf-8") or "utf-8"
        text = raw.decode(encoding, errors="replace").strip()
        if not text:
            return []
        if extension == "md":
            return self._parse_markdown_sections(text)
        return [ParsedBlock(page_content=text, block_type="prose")]

    def _parse_markdown_sections(self, text: str) -> list[ParsedBlock]:
        blocks: list[ParsedBlock] = []
        headings: list[tuple[int, str]] = []
        lines: list[str] = []
        pending_heading = ""

        def heading_path() -> tuple[str, ...]:
            return tuple(title for _, title in headings)

        def flush() -> None:
            content = "\n".join(lines).strip()
            if content:
                blocks.append(ParsedBlock(content, "prose", heading_path=heading_path()))
            lines.clear()

        def flush_standalone_heading() -> None:
            nonlocal pending_heading
            if pending_heading:
                blocks.append(ParsedBlock(
                    pending_heading, "heading", heading_path=heading_path()
                ))
                pending_heading = ""

        for line in text.splitlines():
            match = self._MARKDOWN_HEADING.match(line)
            if match:
                flush()
                flush_standalone_heading()
                level = len(match.group(1))
                title = match.group(2).strip()
                headings[:] = [(heading_level, heading_title)
                               for heading_level, heading_title in headings
                               if heading_level < level]
                headings.append((level, title))
                pending_heading = title
            else:
                lines.append(line)
                if line.strip():
                    pending_heading = ""
        flush()
        flush_standalone_heading()
        return blocks

    @staticmethod
    def _clean_table_cell(value: object) -> str:
        if value is None:
            return ""
        return str(value).strip().replace("\r", " ").replace("\n", " ").replace("|", "\\|")

    @staticmethod
    def _rows_to_markdown(rows: Iterable[list[str]]) -> str:
        normalized_rows = [list(row) for row in rows]
        if not normalized_rows:
            return ""
        width = max(len(row) for row in normalized_rows)
        normalized_rows = [row + [""] * (width - len(row)) for row in normalized_rows]
        header = [value or f"列{index + 1}" for index, value in enumerate(normalized_rows[0])]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]
        lines.extend("| " + " | ".join(row) + " |" for row in normalized_rows[1:])
        return "\n".join(lines)
